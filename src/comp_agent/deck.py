from __future__ import annotations

from pathlib import Path
from dataclasses import asdict, is_dataclass
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt
from PIL import Image

from comp_agent.deck_data import build_comp_study_deck_data
from comp_agent.images import download_hero_images
from comp_agent.models import CompCandidate, CompCriterion, MetricSummary, ProjectBrief, SourceLogEntry
from comp_agent.workspace import write_json


FONT = "Segoe UI"
INK = RGBColor(0, 0, 0)
MUTED = RGBColor(169, 161, 155)
LINE = RGBColor(190, 182, 175)
GRID = RGBColor(169, 161, 155)
BLACK = RGBColor(0, 0, 0)
FILL = RGBColor(255, 255, 255)
PAPER = RGBColor(255, 255, 255)
ACCENT = RGBColor(0, 169, 92)
ACCENT_2 = RGBColor(255, 172, 42)
DARK = RGBColor(0, 0, 0)
LIGHT_GREEN = RGBColor(115, 201, 45)
ACCENT_TINT = RGBColor(229, 246, 237)
MATRIX_DOT_DIAMETER = 0.052
LOGO_PATH = Path(__file__).with_name("assets") / "pelli_clarke_partners_logo.jpg"


def create_concept_deck(
    brief: ProjectBrief,
    criteria: list[CompCriterion],
    candidates: list[CompCandidate],
    metrics: list[MetricSummary],
    source_log: list[SourceLogEntry],
    output_path: str | Path,
    *,
    records: list[Any] | None = None,
    data_output_dir: str | Path | None = None,
) -> Path:
    output = Path(output_path)
    output.parent.mkdir(parents=True, exist_ok=True)
    record_dicts = [dict(record) if isinstance(record, dict) else asdict(record) if is_dataclass(record) else {} for record in (records or [])]
    deck_data = build_comp_study_deck_data(brief, record_dicts, candidates, source_log)
    if data_output_dir:
        download_hero_images(deck_data, Path(data_output_dir) / "images")
    return create_concept_deck_from_data(deck_data, output, data_output_dir=data_output_dir)


def create_concept_deck_from_data(
    deck_data: dict[str, Any],
    output_path: str | Path,
    *,
    data_output_dir: str | Path | None = None,
) -> Path:
    output = Path(output_path)
    output.parent.mkdir(parents=True, exist_ok=True)
    if data_output_dir:
        _write_deck_artifacts(deck_data, Path(data_output_dir) / "json")
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    _cover_slide(prs, deck_data)
    _summary_matrix_slides(prs, deck_data)
    for index, comp in enumerate(deck_data["comps"], start=1):
        _profile_slide(prs, comp, index)
    _comparison_matrix_slides(prs, deck_data)
    _design_levers_slide(prs, deck_data)
    _takeaways_slide(prs, deck_data)

    prs.save(output)
    return output


def _write_deck_artifacts(deck_data: dict[str, Any], data_dir: Path) -> None:
    data_dir.mkdir(parents=True, exist_ok=True)
    write_json(data_dir / "deck_data.json", deck_data)
    write_json(data_dir / "deck_strategy.json", deck_data["deck_strategy"])
    write_json(data_dir / "approved_comps_normalized.json", deck_data["comps"])
    source_metadata = [
        {"project_name": comp["project_name"], "sources": comp["primary_sources"]}
        for comp in deck_data["comps"]
    ]
    write_json(data_dir / "source_metadata.json", source_metadata)
    diligence_notes = [
        {
            "project_name": comp["project_name"],
            "selection_reasoning": comp["selection_reasoning_internal"],
            "diligence_notes": comp["diligence_notes_internal"],
            "data_confidence": comp["data_confidence"],
        }
        for comp in deck_data["comps"]
    ]
    write_json(data_dir / "diligence_notes.json", diligence_notes)


def _blank_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = PAPER
    return slide


def _add_text(
    slide,
    text: Any,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    size: int = 12,
    bold: bool = False,
    color=INK,
    align: PP_ALIGN | None = None,
    valign: MSO_ANCHOR | None = None,
    word_wrap: bool = True,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = word_wrap
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0
    if valign is not None:
        frame.vertical_anchor = valign
    paragraph = frame.paragraphs[0]
    if align is not None:
        paragraph.alignment = align
    paragraph.text = _display(text)
    for paragraph in frame.paragraphs:
        if align is not None:
            paragraph.alignment = align
        runs = paragraph.runs or [paragraph.add_run()]
        for run in runs:
            run.font.name = FONT
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = color
    return box


def _add_section_label(slide, text: str, x: float, y: float, w: float = 4.0) -> None:
    _add_text(slide, text.upper(), x, y, w, 0.18, size=7, bold=True, color=MUTED)


def _add_rule(slide, x: float, y: float, w: float, color: RGBColor = GRID, *, width: float = 0.35) -> None:
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.01))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = color
    shape.line.width = Pt(width)
    _disable_shadow(shape)


def _add_table(slide, rows: list[list[Any]], x: float, y: float, w: float, h: float, *, header_size: int = 8, body_size: int = 7) -> None:
    if not rows:
        return
    table = slide.shapes.add_table(len(rows), len(rows[0]), Inches(x), Inches(y), Inches(w), Inches(h)).table
    for row_index, row in enumerate(rows):
        for col_index, value in enumerate(row):
            cell = table.cell(row_index, col_index)
            cell.text = _display(value)
            cell.fill.solid()
            cell.fill.fore_color.rgb = FILL if row_index == 0 else PAPER
            cell.margin_left = Inches(0.04)
            cell.margin_right = Inches(0.04)
            cell.margin_top = Inches(0.03)
            cell.margin_bottom = Inches(0.03)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.name = FONT
                paragraph.font.size = Pt(header_size if row_index == 0 else body_size)
                paragraph.font.bold = row_index == 0
                paragraph.font.color.rgb = INK if row_index == 0 else DARK


def _cover_slide(prs: Presentation, deck_data: dict[str, Any]) -> None:
    subject = deck_data["subject_project"]
    strategy = deck_data["deck_strategy"]
    slide = _blank_slide(prs)
    context = subject["address"] or ", ".join(part for part in [subject["city"], subject["state_or_country"]] if part)
    intent = strategy.get("cover_intent_label") or strategy.get("project_type_label") or strategy.get("study_focus_label", "")
    _add_text(slide, strategy.get("deck_title", "Comparative Projects"), 0.75, 1.05, 11.0, 0.8, size=42, bold=True)
    _add_text(slide, context, 0.78, 2.05, 10.6, 0.34, size=18, color=ACCENT)
    _add_text(slide, intent, 0.78, 2.55, 8.6, 0.28, size=13, color=MUTED)
    _add_text(slide, f"{len(deck_data['comps'])} approved comps", 0.78, 2.9, 8.2, 0.2, size=9, color=MUTED)
    _add_logo(slide, 0.78, 6.34, 2.6)
    _add_text(slide, f"Generated {_generated_full_date()}", 0.78, 6.72, 3.6, 0.2, size=7, color=MUTED)


def _summary_matrix_slides(prs: Presentation, deck_data: dict[str, Any]) -> None:
    comps = deck_data["comps"]
    columns = deck_data["deck_strategy"]["summary_matrix_columns"]
    for offset in range(0, len(comps), 8):
        slide = _blank_slide(prs)
        _slide_title(slide, "Comp Summary Matrix", "Approved comparable projects")
        rows = [columns]
        for comp in comps[offset : offset + 8]:
            rows.append(
                [
                    comp["project_name"],
                    comp["location"],
                    comp["project_type"],
                    comp["scale"]["display"],
                    comp["status_year"],
                    comp["intervention_type"],
                    comp["relevance_to_subject"],
                ]
            )
        _add_table(slide, rows, 0.55, 1.25, 12.25, 5.75, header_size=8, body_size=6)
        _content_footer(slide, len(prs.slides))


def _profile_slide(prs: Presentation, comp: dict[str, Any], index: int) -> None:
    slide = _blank_slide(prs)
    _add_image_triptych(slide, comp, 0.55, 0.65, 7.2, 5.95)

    _add_text(slide, _profile_title_text(comp["project_name"]), 8.05, 0.62, 5.0, 0.78, size=20, bold=True)
    meta = " | ".join(part for part in [comp["location"], comp["status_year"]] if part and part != "—")
    _add_text(slide, meta, 8.06, 1.48, 4.8, 0.22, size=9, color=MUTED)
    _tags(slide, _profile_tags(comp), 8.05, 1.82, 5.05)
    _facts_card(slide, "Project Profile", _universal_fact_rows(comp), 8.05, 2.16, 5.05, 1.76, row_height=0.19, value_size=6)
    _facts_card(slide, "Study-Specific Insights", list(comp["adaptive_fields"].items())[:5], 8.05, 4.05, 5.05, 1.82, row_height=0.25, value_size=6, max_value_chars=110)
    _add_section_label(slide, "Relevance to Subject", 8.05, 6.03)
    _add_text(slide, comp["relevance_to_subject"], 8.05, 6.25, 5.05, 0.42, size=9)
    _add_text(slide, _sources_footer(comp), 0.55, 6.82, 11.85, 0.18, size=5, color=MUTED)
    _content_footer(slide, len(prs.slides))


def _add_image_triptych(slide, comp: dict[str, Any], x: float, y: float, w: float, h: float) -> None:
    gap = 0.28
    column_w = (w - gap) / 2
    support_h = (h - gap) / 2
    right_x = x + column_w + gap
    _add_package_image(slide, comp, "overall", x, y, column_w, h)
    _add_package_image(slide, comp, "focus", right_x, y, column_w, support_h)
    _add_package_image(slide, comp, "detail", right_x, y + support_h + gap, column_w, support_h)


def _add_package_image(slide, comp: dict[str, Any], slot: str, x: float, y: float, w: float, h: float) -> None:
    package = comp.get("image_package") if isinstance(comp.get("image_package"), dict) else {}
    slot_data = package.get(slot) if isinstance(package.get(slot), dict) else {}
    path = slot_data.get("path") or ((comp.get("hero_image") or {}).get("path") if slot == "overall" else "")
    if path and Path(path).exists() and Path(path).suffix.lower() in {".jpg", ".jpeg", ".png"}:
        try:
            _add_cover_picture(slide, Path(path), x, y, w, h)
            return
        except Exception:
            pass
    image = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    image.fill.solid()
    image.fill.fore_color.rgb = FILL
    image.line.color.rgb = LINE
    _add_text(slide, "Image pending", x + w / 2 - 0.72, y + h / 2, 1.44, 0.2, size=8, color=MUTED)


def _add_cover_picture(slide, path: Path, x: float, y: float, w: float, h: float) -> None:
    with Image.open(path) as image:
        image_w, image_h = image.size
    if image_w <= 0 or image_h <= 0:
        slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w), height=Inches(h))
        return
    image_ratio = image_w / image_h
    frame_ratio = w / h
    picture = slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w), height=Inches(h))
    if image_ratio >= frame_ratio:
        crop = (1 - frame_ratio / image_ratio) / 2
        picture.crop_left = crop
        picture.crop_right = crop
    else:
        crop = (1 - image_ratio / frame_ratio) / 2
        picture.crop_top = crop
        picture.crop_bottom = crop


def _comparison_matrix_slides(prs: Presentation, deck_data: dict[str, Any]) -> None:
    comps = deck_data["comps"]
    columns = deck_data["deck_strategy"]["comparison_matrix_columns"][:18]
    row_capacity = 14
    if not comps:
        _comparison_matrix_slide(prs, deck_data, [], columns, page=1, total_pages=1)
        return
    pages = [comps[index : index + row_capacity] for index in range(0, len(comps), row_capacity)]
    for page_index, page_comps in enumerate(pages, start=1):
        _comparison_matrix_slide(prs, deck_data, page_comps, columns, page=page_index, total_pages=len(pages))


def _comparison_matrix_slide(
    prs: Presentation,
    deck_data: dict[str, Any],
    comps: list[dict[str, Any]],
    columns: list[str],
    *,
    page: int,
    total_pages: int,
) -> None:
    slide = _blank_slide(prs)
    title = _matrix_title(deck_data)
    if total_pages > 1:
        title = f"{title} ({page}/{total_pages})"
    count_label = f"{len(comps)} comps x {len(columns)} features"
    _slide_title(slide, title, count_label)
    if not comps or not columns:
        _add_text(slide, "Comparison data pending.", 0.75, 3.2, 5.0, 0.3, size=14, color=MUTED)
        return

    _matrix_legend(slide, 0.55, 1.25, 11.85)

    left_x = 0.55
    left_w = 2.7
    matrix_x = 3.4
    matrix_w = 8.75
    header_y = 1.48
    header_h = 0.48
    grid_y = 2.08
    row_h = 0.32
    grid_h = row_h * len(comps)
    col_w = matrix_w / len(columns)
    score_w = 0.34
    bottom_y = grid_y + grid_h + 0.22
    bottom_h = 0.24

    _matrix_header(slide, columns, matrix_x, header_y, matrix_w, header_h, col_w)
    _project_name_rail(slide, comps, left_x, grid_y, left_w, row_h)
    _matrix_grid(slide, comps, columns, matrix_x, grid_y, matrix_w, grid_h, col_w, row_h)
    _row_heat_bar(slide, comps, columns, matrix_x + matrix_w + 0.28, grid_y, score_w, row_h)
    _column_heat_bar(slide, columns, comps, matrix_x, bottom_y, matrix_w, bottom_h, col_w)
    _content_footer(slide, len(prs.slides))


def _takeaways_slide(prs: Presentation, deck_data: dict[str, Any]) -> None:
    slide = _blank_slide(prs)
    _slide_title(slide, "Project Positioning Takeaways", "Patterns across the approved comp set")
    summary = deck_data.get("takeaway_summary") or "The comp set points toward a small number of recurring market and design patterns."
    _add_text(slide, summary, 0.55, 1.28, 7.35, 0.55, size=13, color=DARK)
    _takeaway_trend_table(slide, deck_data.get("takeaways") or [], 0.55, 2.08, 11.85, 4.72)
    _content_footer(slide, len(prs.slides))


def _design_levers_slide(prs: Presentation, deck_data: dict[str, Any]) -> None:
    slide = _blank_slide(prs)
    _slide_title(slide, "Design Levers", "Translating comp evidence into PC&P design opportunities")
    _add_text(
        slide,
        "This read groups the market signals into moves the design team can shape, edit, and elevate.",
        0.55,
        1.28,
        8.8,
        0.34,
        size=11,
        color=MUTED,
    )
    groups = _design_lever_groups(deck_data)
    y = 1.82
    row_step = 0.96
    for index, group in enumerate(groups, start=1):
        _add_text(slide, f"{index:02d}", 0.72, y + 0.1, 0.38, 0.18, size=8, bold=True, color=ACCENT)
        _add_text(slide, group["title"], 1.25, y + 0.02, 2.8, 0.22, size=13, bold=True, color=INK)
        _add_text(slide, group["description"], 1.25, y + 0.34, 4.3, 0.24, size=7, color=MUTED)
        _design_lever_tokens(slide, group["features"], len((deck_data.get("comps") or [])), 6.1, y + 0.15)
        if index < len(groups):
            _add_rule(slide, 0.72, y + 0.82, 11.1, LINE)
        y += row_step
    _content_footer(slide, len(prs.slides))


def _matrix_title(deck_data: dict[str, Any]) -> str:
    return "Features and Amenities Matrix"


def _takeaway_trend_table(slide, takeaways: list[dict[str, Any]], x: float, y: float, w: float, h: float) -> None:
    rows = [item.get("trend") or "" for item in takeaways if item.get("trend")]
    rows = rows[:10]
    if not rows:
        rows = ["Comparable projects share a small number of recurring design and positioning moves."]
    header_h = 0.34
    row_h = min(0.44, (h - header_h) / max(1, len(rows)))
    _outline_rect(slide, x, y, w, header_h + row_h * len(rows), BLACK, fill=FILL, line_width=1.0)
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(header_h))
    header.fill.solid()
    header.fill.fore_color.rgb = ACCENT_TINT
    header.line.color.rgb = BLACK
    header.line.width = Pt(0.45)
    _disable_shadow(header)
    _add_text(slide, "Trends Across Comps", x + 0.18, y + 0.09, w - 0.36, 0.14, size=8, bold=True, color=INK)
    for index, trend in enumerate(rows, start=1):
        row_y = y + header_h + (index - 1) * row_h
        if index > 1:
            _add_rule(slide, x, row_y, w)
        _add_text(slide, f"{index:02d}", x + 0.18, row_y, 0.38, row_h, size=7, bold=True, color=ACCENT, valign=MSO_ANCHOR.MIDDLE)
        _add_text(slide, _truncate_for_display(trend, 180), x + 0.7, row_y, w - 0.95, row_h, size=8, color=INK, valign=MSO_ANCHOR.MIDDLE)


def _matrix_legend(slide, x: float, y: float, w: float) -> None:
    _add_text(slide, "Dots indicate presence of a feature; shaded bars show frequency across the comp set.", x, y, 6.5, 0.16, size=7, color=MUTED)
    legend_w = 0.22
    start_x = x + w - 1.55
    for index, ratio in enumerate([0.2, 0.6, 1.0]):
        _outline_rect(slide, start_x + index * (legend_w + 0.03), y - 0.005, legend_w, 0.14, LINE, fill=_heat_color(ratio), line_width=0.25)
    _add_text(slide, "feature frequency", start_x + 0.82, y, 0.78, 0.14, size=5, color=MUTED)


def _design_lever_groups(deck_data: dict[str, Any]) -> list[dict[str, Any]]:
    counts = _matrix_feature_counts(deck_data)
    if not counts:
        return []
    cluster_specs = [
        {
            "title": "Hospitality Arrival",
            "description": "Arrival features are shifting from circulation into hospitality and dwell time.",
            "labels": ["Lobby Seating", "Sky Lobby", "Concierge / Hospitality", "Cafe", "Bar / Lounge"],
        },
        {
            "title": "Food & Social Program",
            "description": "Food, beverage, and lounge uses are recurring tools for daily activity.",
            "labels": ["Cafe", "Restaurant", "Food Hall", "Bar / Lounge", "Tenant Lounge"],
        },
        {
            "title": "Tenant Amenity Network",
            "description": "Shared tenant amenities point to a more complete workday support system.",
            "labels": ["Tenant Lounge", "Co-working", "Conference Center", "Fitness / Wellness", "Event Space", "Bike / Mobility"],
        },
        {
            "title": "Outdoor Amenity Value",
            "description": "Outdoor features show where amenity value extends beyond enclosed space.",
            "labels": ["Terraces", "Public Plaza", "Fitness / Wellness", "Tenant Lounge", "Cafe"],
        },
        {
            "title": "Public-Facing Activation",
            "description": "Street-level features make the project more visible, open, and active.",
            "labels": ["Retail", "Restaurant", "Cafe", "Food Hall", "Public Plaza", "Transit Connection", "Art / Installations"],
        },
        {
            "title": "Access & Mobility",
            "description": "Connection features turn convenience into part of the user experience.",
            "labels": ["Transit Connection", "Bike / Mobility", "Public Plaza", "Retail"],
        },
        {
            "title": "Identity & Place",
            "description": "Identity features help the project read as specific rather than generic.",
            "labels": ["Art / Installations", "Sky Lobby", "Public Plaza", "Concierge / Hospitality", "Food Hall"],
        },
    ]
    groups = []
    for spec in cluster_specs:
        features = _feature_group(counts, spec["labels"])
        if not features:
            continue
        score = sum(count for _label, count in features)
        groups.append({**spec, "features": features, "score": score})
    groups = sorted(groups, key=lambda group: (-group["score"], group["title"]))[:5]
    if not groups:
        groups = [{"title": "Dominant Feature Signals", "description": "The most frequent features define the clearest market signal.", "features": counts[:4], "score": sum(count for _label, count in counts[:4])}]
    return [{"title": group["title"], "description": group["description"], "features": group["features"]} for group in groups]


def _matrix_feature_counts(deck_data: dict[str, Any]) -> list[tuple[str, int]]:
    comps = deck_data.get("comps") or []
    columns = (deck_data.get("deck_strategy") or {}).get("comparison_matrix_columns") or []
    rows: list[tuple[str, int]] = []
    for column in columns:
        count = sum(1 for comp in comps if _has_matrix_dot((comp.get("comparison_flags") or {}).get(column)))
        if count:
            rows.append((column, count))
    return sorted(rows, key=lambda item: (-item[1], item[0]))


def _feature_group(counts: list[tuple[str, int]], labels: list[str]) -> list[tuple[str, int]]:
    lookup = {label: count for label, count in counts}
    selected = [(label, lookup[label]) for label in labels if label in lookup]
    return selected[:4]


def _design_lever_tokens(slide, features: list[tuple[str, int]], comp_count: int, x: float, y: float) -> None:
    cursor = x
    for label, count in features[:4]:
        text = _compact_feature_label(label)
        width = max(0.72, min(1.12, 0.085 * len(text) + 0.34))
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(cursor), Inches(y), Inches(width), Inches(0.24))
        shape.fill.solid()
        shape.fill.fore_color.rgb = ACCENT_TINT
        shape.line.color.rgb = ACCENT
        shape.line.width = Pt(0.45)
        _disable_shadow(shape)
        _add_text(slide, text, cursor + 0.04, y, width - 0.08, 0.24, size=6, bold=True, color=INK, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        _add_text(slide, f"{count}/{max(1, comp_count)}", cursor, y + 0.32, width, 0.12, size=5, bold=True, color=MUTED, align=PP_ALIGN.CENTER)
        cursor += width + 0.18


def _matrix_header(slide, columns: list[str], x: float, y: float, w: float, h: float, col_w: float) -> None:
    _outline_rect(slide, x, y, w, h, BLACK, fill=PAPER, line_width=1.35)
    labels = [_wrap_matrix_label(column, col_w) for column in columns]
    header_font_size = _matrix_header_shared_font_size(labels, col_w)
    for col_index, column in enumerate(columns):
        cx = x + col_index * col_w
        if col_index:
            _add_rule_vertical(slide, cx, y, h, BLACK, width=0.45)
        _add_text(slide, labels[col_index], cx + 0.006, y + 0.04, col_w - 0.012, h - 0.08, size=header_font_size, color=INK, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, word_wrap=False)


def _project_name_rail(slide, comps: list[dict[str, Any]], x: float, y: float, w: float, row_h: float) -> None:
    name_w = w * 0.62
    _outline_rect(slide, x, y, w, row_h * len(comps), BLACK, fill=PAPER, line_width=1.2)
    _add_rule_vertical(slide, x + name_w, y, row_h * len(comps), BLACK, width=0.55)
    for row_index, comp in enumerate(comps):
        cy = y + row_index * row_h
        if row_index:
            _add_rule(slide, x, cy, w, BLACK, width=0.45)
        _add_text(slide, _short_cell(comp.get("project_name"), 32), x + 0.035, cy + 0.045, name_w - 0.06, row_h - 0.04, size=6)
        _add_text(slide, _short_cell(comp.get("location"), 24), x + name_w + 0.035, cy + 0.045, w - name_w - 0.06, row_h - 0.04, size=6)


def _matrix_grid(
    slide,
    comps: list[dict[str, Any]],
    columns: list[str],
    x: float,
    y: float,
    w: float,
    h: float,
    col_w: float,
    row_h: float,
) -> None:
    _outline_rect(slide, x, y, w, h, BLACK, fill=PAPER, line_width=1.35)
    for col_index in range(1, len(columns)):
        _add_rule_vertical(slide, x + col_index * col_w, y, h, GRID, width=0.35)
    for row_index, comp in enumerate(comps):
        cy = y + row_index * row_h
        if row_index:
            _add_rule(slide, x, cy, w, GRID)
        flags = comp.get("comparison_flags") or {}
        for col_index, column in enumerate(columns):
            if _has_matrix_dot(flags.get(column)):
                _add_dot(slide, x + col_index * col_w + col_w / 2, cy + row_h / 2)


def _row_heat_bar(slide, comps: list[dict[str, Any]], columns: list[str], x: float, y: float, w: float, row_h: float) -> None:
    max_count = max(1, len(columns))
    for row_index, comp in enumerate(comps):
        count = sum(1 for column in columns if _has_matrix_dot((comp.get("comparison_flags") or {}).get(column)))
        _outline_rect(slide, x, y + row_index * row_h, w, row_h, BLACK, fill=_heat_color(count / max_count), line_width=0.55)


def _column_heat_bar(slide, columns: list[str], comps: list[dict[str, Any]], x: float, y: float, w: float, h: float, col_w: float) -> None:
    max_count = max(1, len(comps))
    for col_index, column in enumerate(columns):
        count = sum(1 for comp in comps if _has_matrix_dot((comp.get("comparison_flags") or {}).get(column)))
        _outline_rect(slide, x + col_index * col_w, y, col_w, h, BLACK, fill=_heat_color(count / max_count), line_width=0.55)


def _has_matrix_dot(value: Any) -> bool:
    return value not in (None, "", "—", "None", False, 0)


def _add_dot(slide, cx: float, cy: float, diameter: float = MATRIX_DOT_DIAMETER) -> None:
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx - diameter / 2), Inches(cy - diameter / 2), Inches(diameter), Inches(diameter))
    dot.fill.solid()
    dot.fill.fore_color.rgb = RGBColor(0, 0, 0)
    dot.line.color.rgb = RGBColor(0, 0, 0)
    dot.line.width = Pt(0)
    _disable_shadow(dot)


def _heat_color(ratio: float) -> RGBColor:
    ratio = max(0.0, min(1.0, ratio))
    low = (255, 172, 42)
    high = (0, 169, 92)
    values = [round(low[index] + (high[index] - low[index]) * ratio) for index in range(3)]
    return RGBColor(*values)


def _outline_rect(slide, x: float, y: float, w: float, h: float, line: RGBColor, *, fill: RGBColor = PAPER, line_width: float = 0.5) -> None:
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(line_width)
    _disable_shadow(shape)


def _add_rule_vertical(slide, x: float, y: float, h: float, color: RGBColor = LINE, *, width: float = 0.35) -> None:
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(0.006), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = color
    shape.line.width = Pt(width)
    _disable_shadow(shape)


def _disable_shadow(shape) -> None:
    try:
        shape.shadow.inherit = False
    except Exception:
        pass


def _wrap_matrix_label(value: Any, col_w: float = 0.7) -> str:
    compact = _compact_feature_label(value)
    stacked = {
        "Lobby Seating": "Lobby\nSeating",
        "Sky Lobby": "Sky\nLobbies",
        "Food Hall": "Food\nHalls",
        "Public Plaza": "Public\nPlazas",
        "Tenant Lounge": "Tenant\nLounges",
        "Bar / Lounge": "Bars",
        "Bike / Mobility": "Bike\nRooms",
    }
    return stacked.get(str(value or ""), compact)


def _compact_feature_label(value: Any) -> str:
    text = " ".join(str(value or "").replace(" / ", "/").replace("/", " ").replace("-", "").split())
    labels = {
        "Art Installations": "Art",
        "Bar Lounge": "Bar",
        "Bike Mobility": "Bike",
        "Branded Amenities": "Brand",
        "Concierge Hospitality": "Concierge",
        "Conference Center": "Meetings",
        "Co working": "Cowork",
        "Coworking": "Cowork",
        "Fitness Wellness": "Wellness",
        "Food Hall": "Food Halls",
        "Lobby Seating": "Seating",
        "Public Plaza": "Plaza",
        "Public Realm": "Realm",
        "Sky Lobby": "Sky Lobbies",
        "Tenant Lounge": "Lounges",
        "Transit Connection": "Transit",
        "Event Space": "Events",
        "Sustainability": "Green",
        "Technology": "Tech",
    }
    return labels.get(text, text)


def _matrix_header_font_size(label: str, col_w: float) -> int:
    longest = max((len(line) for line in label.splitlines()), default=0)
    line_count = len(label.splitlines())
    if line_count >= 2:
        return 5 if col_w < 0.7 or longest > 7 else 6
    if col_w < 0.64 or longest > 8:
        return 5
    if col_w < 0.78 or longest > 6:
        return 6
    return 7


def _matrix_header_shared_font_size(labels: list[str], col_w: float) -> int:
    if not labels:
        return 7
    return min(_matrix_header_font_size(label, col_w) for label in labels)


def _trim_matrix_label_line(line: str, limit: int = 14) -> str:
    if len(line) <= limit:
        return line
    return line[:limit].rsplit(" ", 1)[0].strip() or line[:limit]


def _short_cell(value: Any, limit: int) -> str:
    text = " ".join(_display(value).split())
    if len(text) <= limit:
        return text
    return text[:limit].rsplit(" ", 1)[0].strip() or text[:limit]


def _slide_title(slide, title: str, subtitle: str) -> None:
    _add_text(slide, title, 0.55, 0.38, 8.8, 0.42, size=22, bold=True)
    _add_text(slide, subtitle, 0.57, 0.82, 8.8, 0.2, size=9, color=MUTED)
    _add_rule(slide, 0.55, 1.08, 12.25)


def _truncate_for_display(text: Any, max_chars: int) -> str:
    s = "" if text is None else str(text)
    if len(s) <= max_chars:
        return s
    return s[: max(1, max_chars - 1)].rstrip() + "…"


def _facts_card(
    slide,
    title: str,
    rows: list[tuple[Any, Any]],
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    row_height: float = 0.21,
    value_size: int = 7,
    max_value_chars: int | None = None,
) -> None:
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = FILL
    shape.line.color.rgb = LINE
    _add_section_label(slide, title, x + 0.18, y + 0.15, w - 0.35)
    row_y = y + 0.43
    bottom = y + h - 0.13
    for label, value in rows:
        if row_y + row_height > bottom:
            break
        display_value = _truncate_for_display(value, max_value_chars) if max_value_chars else value
        _add_text(slide, label, x + 0.18, row_y, 1.25, 0.18, size=5, bold=True, color=MUTED)
        _add_text(slide, display_value, x + 1.47, row_y, w - 1.65, row_height, size=value_size)
        row_y += row_height


def _tags(slide, tags: list[str], x: float, y: float, w: float = 4.65) -> None:
    tags = tags[:5]
    if not tags:
        return
    gap = 0.08
    available = w - gap * (len(tags) - 1)
    weights = [max(0.8, min(2.2, 0.08 * len(tag) + 0.45)) for tag in tags]
    total_weight = sum(weights)
    cursor = x
    for index, tag in enumerate(tags):
        width = available * weights[index] / total_weight
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(cursor), Inches(y), Inches(width), Inches(0.24))
        shape.fill.solid()
        shape.fill.fore_color.rgb = ACCENT_TINT
        shape.line.color.rgb = ACCENT
        _add_text(slide, tag, cursor + 0.04, y, width - 0.08, 0.24, size=5, bold=True, color=INK, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        cursor += width + gap


def _profile_title_text(value: Any) -> str:
    title = " ".join(_display(value).split())
    if len(title) <= 34:
        return title
    words = title.split()
    midpoint = len(title) / 2
    best_index = 1
    best_delta = len(title)
    cursor = 0
    for index, word in enumerate(words[:-1], start=1):
        cursor += len(word) + (1 if index > 1 else 0)
        delta = abs(cursor - midpoint)
        if delta < best_delta:
            best_delta = delta
            best_index = index
    lines = [" ".join(words[:best_index]), " ".join(words[best_index:])]
    return "\n".join(_fit_title_line(line) for line in lines[:2])


def _fit_title_line(line: str) -> str:
    if len(line) <= 34:
        return line
    return line[:34].rsplit(" ", 1)[0].strip() or line[:34]


def _profile_tags(comp: dict[str, Any]) -> list[str]:
    values = [key for key, value in (comp.get("comparison_flags") or {}).items() if _has_matrix_dot(value)]
    if len(values) < 3:
        values.extend(key for key, value in (comp.get("adaptive_fields") or {}).items() if value)
    if len(values) < 3:
        values.extend([comp["key_program"], comp["intervention_type"], comp["project_type"]])
    tags = []
    for value in values:
        short = _compact_feature_label(value)
        if short and short != "—" and short not in tags:
            tags.append(short)
    return tags


def _universal_fact_rows(comp: dict[str, Any]) -> list[tuple[str, str]]:
    return [
        ("Program", comp["project_type"]),
        ("Scale", comp["scale"]["display"]),
        ("Year / Status", comp["status_year"]),
        ("Owner / Developer", comp["owner_developer"]),
        ("Architect / Designer", comp["architect_designer"]),
        ("Intervention", comp["intervention_type"]),
    ]


def _sources_footer(comp: dict[str, Any]) -> str:
    sources = comp.get("primary_sources") or []
    if not sources:
        return "Sources: pending verification"
    labels = []
    for source in sources[:4]:
        label = source.get("publisher") or source.get("title") or source.get("source_type")
        if label and label not in labels:
            labels.append(label)
    return "Sources: " + "; ".join(labels)


def _display(value: Any) -> str:
    if value in (None, "", "None"):
        return "—"
    return str(value)


def _generated_full_date() -> str:
    from datetime import datetime

    now = datetime.now()
    return f"{now.strftime('%B')} {now.day}, {now.year}"


def _add_logo(slide, x: float, y: float, w: float) -> None:
    if not LOGO_PATH.exists():
        return
    try:
        with Image.open(LOGO_PATH) as image:
            image_w, image_h = image.size
        h = w * image_h / image_w if image_w else 0.18
        slide.shapes.add_picture(str(LOGO_PATH), Inches(x), Inches(y), width=Inches(w), height=Inches(h))
    except Exception:
        return


def _content_footer(slide, page_number: int) -> None:
    footer_y = 7.13
    logo_w = 1.8
    logo_h = _logo_height(logo_w)
    _add_logo(slide, 0.55, footer_y, logo_w)
    _add_text(slide, f"{page_number:02d}", 12.35, footer_y + logo_h - 0.14, 0.35, 0.14, size=6, color=MUTED, align=PP_ALIGN.RIGHT, valign=MSO_ANCHOR.BOTTOM)


def _logo_height(w: float) -> float:
    if not LOGO_PATH.exists():
        return 0.14
    try:
        with Image.open(LOGO_PATH) as image:
            image_w, image_h = image.size
        return w * image_h / image_w if image_w else 0.14
    except Exception:
        return 0.14
